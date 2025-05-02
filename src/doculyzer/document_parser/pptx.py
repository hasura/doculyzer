"""
PPTX document parser module for the document pointer system.

This module parses PowerPoint (PPTX) documents into structured elements.
"""

import json
import logging
import os
import re
from typing import Dict, Any, List, Optional, Union

try:
    import pptx
    # noinspection PyUnresolvedReferences
    from pptx import Presentation
    from pptx.shapes.autoshape import Shape
    # noinspection PyUnresolvedReferences
    from pptx.shapes.group import GroupShape
    # noinspection PyUnresolvedReferences
    from pptx.shapes.picture import Picture
    # noinspection PyUnresolvedReferences
    from pptx.slide import Slide, SlideLayout
    # noinspection PyUnresolvedReferences
    from pptx.text.text import TextFrame

    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logging.warning("python-pptx not available. Install with 'pip install python-pptx' to use PPTX parser")

from .base import DocumentParser

logger = logging.getLogger(__name__)


class PptxParser(DocumentParser):
    """Parser for PowerPoint (PPTX) documents."""

    def __init__(self, config: Optional[Dict[str, Any]] = None):
        """Initialize the PPTX parser."""
        super().__init__(config)

        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx is required for PPTX parsing")

        # Configuration options
        self.config = config or {}
        self.extract_notes = self.config.get("extract_notes", True)
        self.extract_hidden_slides = self.config.get("extract_hidden_slides", False)
        self.extract_comments = self.config.get("extract_comments", True)
        self.extract_shapes = self.config.get("extract_shapes", True)
        self.extract_images = self.config.get("extract_images", True)
        self.extract_tables = self.config.get("extract_tables", True)
        self.extract_text_boxes = self.config.get("extract_text_boxes", True)
        self.extract_charts = self.config.get("extract_charts", True)
        self.extract_masters = self.config.get("extract_masters", False)
        self.extract_templates = self.config.get("extract_templates", False)
        self.temp_dir = self.config.get("temp_dir", os.path.join(os.path.dirname(__file__), 'temp'))

    def _resolve_element_text(self, location_data: Dict[str, Any], source_content: Optional[Union[str, bytes]]) -> str:
        """
        Resolve the plain text representation of a PPTX element.

        Args:
            location_data: Content location data
            source_content: Optional preloaded source content

        Returns:
            Plain text representation of the element
        """
        # Get the content using the improved _resolve_element_content method
        content = self._resolve_element_content(location_data, source_content)
        element_type = location_data.get("type", "")

        # Handle specific element types
        if element_type == "presentation_body":
            return content.strip()

        elif element_type == "slide":
            return content.strip()

        elif element_type == "text_box" or element_type == "paragraph":
            return content.strip()

        elif element_type == "table" or element_type == "table_cell":
            # The improved _resolve_element_content already formats tables properly
            return content.strip()

        elif element_type == "slide_notes":
            return content.strip()

        elif element_type == "comment":
            # For comments, extract just the comment text without metadata
            if ": " in content:
                return content.split(": ", 1)[1].strip()
            return content.strip()

        elif element_type == "image":
            if "Alt text: " in content:
                return content.split("Alt text: ", 1)[1].strip()
            return "Image"

        elif element_type == "chart":
            if "Chart: " in content and "\n" in content:
                return content.split("\n")[0].replace("Chart: ", "").strip()
            return content.strip()

        # Default
        return content.strip()

    def _resolve_element_content(self, location_data: Dict[str, Any],
                                 source_content: Optional[Union[str, bytes]] = None) -> str:
        """
        Resolve content for specific PPTX element types.

        Args:
            location_data: Content location data
            source_content: Optional pre-loaded source content

        Returns:
            Resolved content string
        """
        source = location_data.get("source", "")
        element_type = location_data.get("type", "")
        slide_index = location_data.get("slide_index", 0)

        # Load the document if source content is not provided
        presentation = None
        temp_file = None
        try:
            if source_content is None:
                # Check if source is a file path
                if os.path.exists(source):
                    try:
                        presentation = Presentation(source)
                    except Exception as e:
                        raise ValueError(f"Error loading PPTX document: {str(e)}")
                else:
                    raise ValueError(f"Source file not found: {source}")
            else:
                # Save content to a temporary file
                if not os.path.exists(self.temp_dir):
                    os.makedirs(self.temp_dir, exist_ok=True)

                import uuid
                temp_file = os.path.join(self.temp_dir, f"temp_{uuid.uuid4().hex}.pptx")
                with open(temp_file, 'wb') as f:
                    if isinstance(source_content, str):
                        f.write(source_content.encode('utf-8'))
                    else:
                        f.write(source_content)

                # Load the document
                try:
                    presentation = Presentation(temp_file)
                except Exception as e:
                    raise ValueError(f"Error loading PPTX document: {str(e)}")

            # Handle different element types
            if element_type == "presentation_body":
                # Return basic presentation information
                slide_count = len(presentation.slides)
                return f"Presentation with {slide_count} slides"

            elif element_type == "slide":
                # Check if slide index is valid
                if slide_index < 0 or slide_index >= len(presentation.slides):
                    return f"Invalid slide index: {slide_index}. Presentation has {len(presentation.slides)} slides."

                # Get the slide
                slide = presentation.slides[slide_index]

                # Extract all text from the slide
                all_text = []
                for shape in slide.shapes:
                    if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                        text = shape.text_frame.text
                        if text:
                            all_text.append(text)

                return "\n\n".join(all_text)

            elif element_type == "text_box":
                # Extract text from a text box shape
                shape_path = location_data.get("shape_path", "")
                if not shape_path:
                    return "No shape path specified"

                # Check if slide index is valid
                if slide_index < 0 or slide_index >= len(presentation.slides):
                    return f"Invalid slide index: {slide_index}"

                # Get the slide
                slide = presentation.slides[slide_index]

                # Parse shape path to locate the shape
                shape_indices = shape_path.split('/')
                shape = self._find_shape_by_path(slide.shapes, shape_indices)

                if not shape or not hasattr(shape, 'has_text_frame') or not shape.has_text_frame:
                    return f"Text shape not found at path: {shape_path}"

                return shape.text_frame.text

            elif element_type == "paragraph":
                # Extract specific paragraph from a text shape
                shape_path = location_data.get("shape_path", "")
                paragraph_index = location_data.get("paragraph_index", 0)

                # Check if slide index is valid
                if slide_index < 0 or slide_index >= len(presentation.slides):
                    return f"Invalid slide index: {slide_index}"

                # Get the slide
                slide = presentation.slides[slide_index]

                # Parse shape path to locate the shape
                shape_indices = shape_path.split('/')
                shape = self._find_shape_by_path(slide.shapes, shape_indices)

                if not shape or not hasattr(shape, 'has_text_frame') or not shape.has_text_frame:
                    return f"Text shape not found at path: {shape_path}"

                # Check if paragraph index is valid
                if not hasattr(shape.text_frame, 'paragraphs') or paragraph_index >= len(shape.text_frame.paragraphs):
                    return f"Invalid paragraph index: {paragraph_index}"

                return shape.text_frame.paragraphs[paragraph_index].text

            elif element_type == "table":
                # Extract table content with proper formatting
                shape_path = location_data.get("shape_path", "")

                # Check if slide index is valid
                if slide_index < 0 or slide_index >= len(presentation.slides):
                    return f"Invalid slide index: {slide_index}"

                # Get the slide
                slide = presentation.slides[slide_index]

                # Parse shape path to locate the shape
                shape_indices = shape_path.split('/')
                shape = self._find_shape_by_path(slide.shapes, shape_indices)

                if not shape or not hasattr(shape, 'has_table') or not shape.has_table:
                    return f"Table shape not found at path: {shape_path}"

                # Get table object
                table = shape.table

                # Format table with consistent structure
                rows_text = []

                # Process each row
                for row in table.rows:
                    cells_text = []
                    for cell in row.cells:
                        cell_text = cell.text_frame.text.strip() if hasattr(cell, 'text_frame') else ""
                        cells_text.append(cell_text)

                    # Join cells with pipe separator
                    rows_text.append(" | ".join(cells_text))

                # Return formatted table
                return "\n".join(rows_text)

            elif element_type == "table_cell":
                # Extract cell content from a table
                shape_path = location_data.get("shape_path", "")
                row = location_data.get("row", 0)
                col = location_data.get("col", 0)

                # Check if slide index is valid
                if slide_index < 0 or slide_index >= len(presentation.slides):
                    return f"Invalid slide index: {slide_index}"

                # Get the slide
                slide = presentation.slides[slide_index]

                # Parse shape path to locate the shape
                shape_indices = shape_path.split('/')
                shape = self._find_shape_by_path(slide.shapes, shape_indices)

                if not shape or not hasattr(shape, 'has_table') or not shape.has_table:
                    return f"Table shape not found at path: {shape_path}"

                # Get table object
                table = shape.table

                # Check if row and column indices are valid
                if row < 0 or row >= len(table.rows) or col < 0 or col >= len(table.columns):
                    return f"Invalid cell coordinates: row={row}, col={col}"

                # Get cell text
                cell = table.cell(row, col)
                return cell.text_frame.text if hasattr(cell, 'text_frame') else ""

            elif element_type == "slide_notes":
                # Extract notes from a slide
                # Check if slide index is valid
                if slide_index < 0 or slide_index >= len(presentation.slides):
                    return f"Invalid slide index: {slide_index}"

                # Get the slide
                slide = presentation.slides[slide_index]

                # Check if slide has notes
                if not hasattr(slide, 'notes_slide') or not slide.notes_slide:
                    return "No notes for this slide"

                # Return notes text
                return slide.notes_slide.notes_text_frame.text

            elif element_type == "comment":
                # Extract a specific comment
                comment_index = location_data.get("comment_index", 0)

                # Check if slide index is valid
                if slide_index < 0 or slide_index >= len(presentation.slides):
                    return f"Invalid slide index: {slide_index}"

                # Get the slide
                slide = presentation.slides[slide_index]

                # Check if slide has comments
                if not hasattr(slide, 'comments') or not slide.comments:
                    return "No comments for this slide"

                # Check if comment index is valid
                if comment_index < 0 or comment_index >= len(slide.comments):
                    return f"Invalid comment index: {comment_index}"

                # Get comment
                comment = slide.comments[comment_index]

                # Format comment details
                author = comment.author if hasattr(comment, 'author') else "Unknown"
                text = comment.text if hasattr(comment, 'text') else ""
                date = comment.date if hasattr(comment, 'date') else None

                if date:
                    return f"Comment by {author} on {date}: {text}"
                else:
                    return f"Comment by {author}: {text}"

            elif element_type == "image":
                # Return information about an image
                shape_path = location_data.get("shape_path", "")

                # Check if slide index is valid
                if slide_index < 0 or slide_index >= len(presentation.slides):
                    return f"Invalid slide index: {slide_index}"

                # Get the slide
                slide = presentation.slides[slide_index]

                # Parse shape path to locate the shape
                shape_indices = shape_path.split('/')
                shape = self._find_shape_by_path(slide.shapes, shape_indices)

                if not shape or not isinstance(shape, Picture):
                    return f"Image shape not found at path: {shape_path}"

                # Get image details
                image_name = shape.image.filename if hasattr(shape, 'image') and hasattr(shape.image,
                                                                                         'filename') else "Unknown"
                alt_text = shape.alt_text if hasattr(shape, 'alt_text') else ""

                return f"Image: {image_name}\nAlt text: {alt_text}"

            elif element_type == "chart":
                # Return information about a chart
                shape_path = location_data.get("shape_path", "")

                # Check if slide index is valid
                if slide_index < 0 or slide_index >= len(presentation.slides):
                    return f"Invalid slide index: {slide_index}"

                # Get the slide
                slide = presentation.slides[slide_index]

                # Parse shape path to locate the shape
                shape_indices = shape_path.split('/')
                shape = self._find_shape_by_path(slide.shapes, shape_indices)

                if not shape or not hasattr(shape, 'has_chart') or not shape.has_chart:
                    return f"Chart shape not found at path: {shape_path}"

                # Get chart details
                chart = shape.chart
                chart_type = str(chart.chart_type) if hasattr(chart, 'chart_type') else "Unknown"
                chart_title = chart.chart_title.text_frame.text if hasattr(chart, 'chart_title') and hasattr(
                    chart.chart_title, 'text_frame') else "Untitled Chart"

                # Get categories and series if available
                categories = []
                series_names = []

                if hasattr(chart, 'plots') and chart.plots:
                    plot = chart.plots[0]

                    if hasattr(plot, 'categories'):
                        for category in plot.categories:
                            if category:
                                categories.append(str(category))

                    if hasattr(plot, 'series'):
                        for series in plot.series:
                            if hasattr(series, 'name') and series.name:
                                series_names.append(str(series.name))

                # Format chart description
                description = f"Chart: {chart_title}\nType: {chart_type}"

                if categories:
                    description += f"\nCategories: {', '.join(categories)}"

                if series_names:
                    description += f"\nSeries: {', '.join(series_names)}"

                return description

            else:
                # For other element types or if no specific handler,
                # return basic information about the presentation
                return f"PowerPoint presentation with {len(presentation.slides)} slides"

        finally:
            # Clean up temporary file
            if temp_file and os.path.exists(temp_file):
                try:
                    os.remove(temp_file)
                except Exception as e:
                    logger.warning(f"Failed to delete temporary file {temp_file}: {str(e)}")

    def supports_location(self, content_location: str) -> bool:
        """
        Check if this parser supports resolving the given location.

        Args:
            content_location: Content location pointer

        Returns:
            True if supported, False otherwise
        """
        try:
            location_data = json.loads(content_location)
            source = location_data.get("source", "")

            # Check if source exists and is a file
            if not os.path.exists(source) or not os.path.isfile(source):
                return False

            # Check file extension for PPTX
            _, ext = os.path.splitext(source.lower())
            return ext in ['.pptx', '.pptm']

        except (json.JSONDecodeError, TypeError):
            return False

    def _extract_links(self, content: str, element_id: str) -> List[Dict[str, Any]]:
        """
        Base method for extracting links from content.

        Args:
            content: Text content
            element_id: ID of the element containing the links

        Returns:
            List of extracted link dictionaries
        """
        links = []

        # Extract URLs using a regular expression
        url_pattern = r'https?://[^\s<>)"\']+|www\.[^\s<>)"\']+|ftp://[^\s<>)"\']+|file://[^\s<>)"\']+|mailto:[^\s<>)"\']+|[^\s<>)"\']+\.(?:com|org|net|edu|gov|io|ai|app)[^\s<>)"\']*'
        urls = re.findall(url_pattern, content)

        # Create link entries for each URL found
        for url in urls:
            # Clean up URL
            if url.startswith('www.'):
                url = 'http://' + url

            links.append({
                "source_id": element_id,
                "link_text": url,
                "link_target": url,
                "link_type": "url"
            })

        # Look for slide references (e.g., "See slide 5")
        slide_refs = re.findall(r'slide\s+(\d+)', content, re.IGNORECASE)

        for ref in slide_refs:
            try:
                slide_num = int(ref)

                links.append({
                    "source_id": element_id,
                    "link_text": f"Slide {slide_num}",
                    "link_target": f"slide_{slide_num}",
                    "link_type": "slide_reference"
                })
            except ValueError:
                pass

        return links

    @staticmethod
    def _extract_document_links(presentation, elements: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
        """
        Helper method to extract hyperlinks from PowerPoint presentation.
        This is called during the parsing phase.

        Args:
            presentation: The PPTX presentation
            elements: List of extracted elements

        Returns:
            List of hyperlink dictionaries
        """
        links = []

        # Map element IDs to elements for quick lookup
        # element_map = {elem["element_id"]: elem for elem in elements}

        # Extract hyperlinks from text shapes
        for element in elements:
            if element["element_type"] in ["text_box", "paragraph", "table_cell"]:
                element_id = element["element_id"]
                text = element.get("metadata", {}).get("text", "")

                # Look for hyperlink patterns in text
                url_pattern = r'https?://[^\s<>)"\']+|www\.[^\s<>)"\']+|ftp://[^\s<>)"\']+|file://[^\s<>)"\']+|mailto:[^\s<>)"\']+|[^\s<>)"\']+\.(?:com|org|net|edu|gov|io|ai|app)[^\s<>)"\']*'
                urls = re.findall(url_pattern, text)

                for url in urls:
                    # Clean up URL
                    if url.startswith('www.'):
                        url = 'http://' + url

                    # Add link
                    links.append({
                        "source_id": element_id,
                        "link_text": url,
                        "link_target": url,
                        "link_type": "url"
                    })

                # Look for slide references (e.g., "See slide 5")
                slide_refs = re.findall(r'slide\s+(\d+)', text, re.IGNORECASE)

                for ref in slide_refs:
                    try:
                        slide_num = int(ref)
                        # Adjust for 0-based indexing
                        # slide_idx = slide_num - 1

                        # Find target slide element
                        target_slide = None
                        for slide_elem in elements:
                            if (slide_elem["element_type"] == "slide" and
                                    slide_elem.get("metadata", {}).get("number") == slide_num):
                                target_slide = slide_elem
                                break

                        if target_slide:
                            # Add link
                            links.append({
                                "source_id": element_id,
                                "link_text": f"Slide {slide_num}",
                                "link_target": target_slide["element_id"],
                                "link_type": "slide_reference"
                            })
                    except (ValueError, IndexError):
                        pass

        return links

    def parse(self, doc_content: Dict[str, Any]) -> Dict[str, Any]:
        """
        Parse a PPTX document into structured elements.

        Args:
            doc_content: Document content and metadata

        Returns:
            Dictionary with document metadata, elements, relationships, and extracted links
        """
        # Extract metadata from doc_content
        source_id = doc_content["id"]
        metadata = doc_content.get("metadata", {}).copy()  # Make a copy to avoid modifying original

        # Check if we have a binary path or content
        binary_path = doc_content.get("binary_path")
        binary_content = doc_content.get("content")

        # If we have content but no path, save it to a temp file
        if not binary_path and binary_content:
            if not os.path.exists(self.temp_dir):
                os.makedirs(self.temp_dir, exist_ok=True)

            import uuid
            temp_file_path = os.path.join(self.temp_dir, f"temp_{uuid.uuid4().hex}.pptx")
            with open(temp_file_path, 'wb') as f:
                if isinstance(binary_content, str):
                    f.write(binary_content.encode('utf-8'))
                else:
                    f.write(binary_content)

            binary_path = temp_file_path

        if not binary_path:
            raise ValueError("PPTX parser requires either binary_path or content to process the presentation")

        # Generate document ID if not present
        doc_id = metadata.get("doc_id", self._generate_id("doc_"))

        # Load PPTX document
        try:
            presentation = Presentation(binary_path)
        except Exception as e:
            logger.error(f"Error loading PPTX document: {str(e)}")
            raise

        # Create document record with metadata
        document = {
            "doc_id": doc_id,
            "doc_type": "pptx",
            "source": source_id,
            "metadata": self._extract_document_metadata(presentation, metadata),
            "content_hash": doc_content.get("content_hash", "")
        }

        # Create root element
        elements = [self._create_root_element(doc_id, source_id)]
        root_id = elements[0]["element_id"]

        # Parse document elements
        elements.extend(self._parse_presentation(presentation, doc_id, root_id, source_id))

        # Extract links from the document using the helper method
        links = self._extract_document_links(presentation, elements)

        # Clean up temporary file if we created one
        if binary_path != doc_content.get("binary_path") and os.path.exists(binary_path):
            try:
                os.remove(binary_path)
                logger.debug(f"Deleted temporary file: {binary_path}")
            except Exception as e:
                logger.warning(f"Failed to delete temporary file {binary_path}: {str(e)}")

        # Return the parsed document with extracted links
        return {
            "document": document,
            "elements": elements,
            "links": links,
            "relationships": []
        }

    @staticmethod
    def _extract_document_metadata(presentation: Presentation, base_metadata: Dict[str, Any]) -> Dict[str, Any]:
        """
        Extract metadata from PPTX document.

        Args:
            presentation: The PPTX presentation
            base_metadata: Base metadata from content source

        Returns:
            Dictionary of document metadata
        """
        # Combine base metadata with document properties
        metadata = base_metadata.copy()

        try:
            # Add core properties
            core_props = presentation.core_properties
            if core_props:
                # Add core properties to metadata
                if core_props.title:
                    metadata["title"] = core_props.title
                if core_props.author:
                    metadata["author"] = core_props.author
                if core_props.created:
                    metadata["created"] = core_props.created.timestamp() if hasattr(core_props.created,
                                                                                    'timestamp') else str(
                        core_props.created)
                if core_props.modified:
                    metadata["modified"] = core_props.modified.timestamp() if hasattr(core_props.modified,
                                                                                      'timestamp') else str(
                        core_props.modified)
                if core_props.last_modified_by:
                    metadata["last_modified_by"] = core_props.last_modified_by
                if core_props.keywords:
                    metadata["keywords"] = core_props.keywords
                if core_props.subject:
                    metadata["subject"] = core_props.subject
                if core_props.comments:
                    metadata["comments"] = core_props.comments
                if core_props.category:
                    metadata["category"] = core_props.category

            # Add presentation statistics
            metadata["slide_count"] = len(presentation.slides)
            metadata["slide_width"] = presentation.slide_width
            metadata["slide_height"] = presentation.slide_height

        except Exception as e:
            logger.warning(f"Error extracting document metadata: {str(e)}")

        return metadata

    def _parse_presentation(self, presentation: Presentation, doc_id: str, parent_id: str, source_id: str) -> List[
        Dict[str, Any]]:
        """
        Parse PowerPoint presentation into structured elements.

        Args:
            presentation: The PPTX presentation
            doc_id: Document ID
            parent_id: Parent element ID
            source_id: Source identifier

        Returns:
            List of parsed elements
        """
        elements: List = []

        # Create presentation body element
        body_id = self._generate_id("body_")
        body_element = {
            "element_id": body_id,
            "doc_id": doc_id,
            "element_type": "presentation_body",
            "parent_id": parent_id,
            "content_preview": "Presentation body",
            "content_location": json.dumps({
                "source": source_id,
                "type": "presentation_body"
            }),
            "content_hash": "",
            "metadata": {
                "slide_count": len(presentation.slides)
            }
        }
        elements.append(body_element)

        # Process slides
        for slide_idx, slide in enumerate(presentation.slides):
            # Skip hidden slides if not configured to extract them
            if not self.extract_hidden_slides and hasattr(slide, 'hidden') and slide.hidden:
                continue

            # Process this slide
            slide_elements = self._process_slide(slide, slide_idx, doc_id, body_id, source_id)
            elements.extend(slide_elements)

        # Extract masters if configured
        if self.extract_masters and hasattr(presentation, 'slide_masters'):
            master_elements = self._extract_slide_masters(presentation, doc_id, parent_id, source_id)
            elements.extend(master_elements)

        # Extract templates if configured
        if self.extract_templates and hasattr(presentation, 'slide_layouts'):
            template_elements = self._extract_slide_templates(presentation, doc_id, parent_id, source_id)
            elements.extend(template_elements)

        return elements

    def _find_shape_by_path(self, shapes, shape_indices):
        """
        Find a shape by following the shape path indices.

        Args:
            shapes: Collection of shapes to search in
            shape_indices: List of indices to follow

        Returns:
            The shape if found, None otherwise
        """
        try:
            if not shape_indices:
                return None

            # Get the first index
            current_idx = int(shape_indices[0])

            # Check if index is valid
            if current_idx < 0 or current_idx >= len(shapes):
                return None

            # Get the shape at this index
            shape = shapes[current_idx]

            # If this is the last index, return the shape
            if len(shape_indices) == 1:
                return shape

            # If this is a group shape, recurse into it
            if isinstance(shape, GroupShape) and hasattr(shape, 'shapes'):
                return self._find_shape_by_path(shape.shapes, shape_indices[1:])

            # If we get here, the path is invalid
            return None

        except (ValueError, IndexError, TypeError):
            return None

    def _process_slide(self, slide: Slide, slide_idx: int, doc_id: str, parent_id: str, source_id: str) -> List[
        Dict[str, Any]]:
        """
        Process a PowerPoint slide into structured elements.

        Args:
            slide: The PPTX slide
            slide_idx: Slide index (0-based)
            doc_id: Document ID
            parent_id: Parent element ID
            source_id: Source identifier

        Returns:
            List of slide-related elements
        """
        elements = []

        # Generate slide ID
        slide_id = self._generate_id(f"slide_{slide_idx}_")

        # Get slide title if available
        slide_title = self._get_slide_title(slide)

        # Create slide element
        slide_element = {
            "element_id": slide_id,
            "doc_id": doc_id,
            "element_type": "slide",
            "parent_id": parent_id,
            "content_preview": f"Slide {slide_idx + 1}: {slide_title}" if slide_title else f"Slide {slide_idx + 1}",
            "content_location": json.dumps({
                "source": source_id,
                "type": "slide",
                "slide_index": slide_idx
            }),
            "content_hash": self._generate_hash(f"slide_{slide_idx}"),
            "metadata": {
                "index": slide_idx,
                "number": slide_idx + 1,
                "title": slide_title,
                "layout": self._get_slide_layout_name(slide),
                "has_notes": bool(slide.notes_slide and slide.notes_slide.notes_text_frame.text),
                "shape_count": len(slide.shapes)
            }
        }
        elements.append(slide_element)

        # Process slide shapes
        if self.extract_shapes:
            shape_elements = self._process_shapes(slide.shapes, doc_id, slide_id, source_id, slide_idx)
            elements.extend(shape_elements)

        # Process slide notes
        if self.extract_notes and slide.notes_slide and slide.notes_slide.notes_text_frame.text:
            notes_elements = self._process_notes(slide.notes_slide, slide_idx, doc_id, slide_id, source_id)
            elements.extend(notes_elements)

        # Process slide comments if available and configured
        if self.extract_comments and hasattr(slide, 'comments'):
            comment_elements = self._process_comments(slide, slide_idx, doc_id, slide_id, source_id)
            elements.extend(comment_elements)

        return elements

    def _process_shapes(self, shapes, doc_id: str, parent_id: str, source_id: str, slide_idx: int,
                        shape_path: str = "") -> List[Dict[str, Any]]:
        """
        Process PowerPoint shapes into structured elements.

        Args:
            shapes: Collection of shapes
            doc_id: Document ID
            parent_id: Parent element ID
            source_id: Source identifier
            slide_idx: Slide index
            shape_path: Path to the shape (for nested shapes)

        Returns:
            List of shape-related elements
        """
        elements = []

        for shape_idx, shape in enumerate(shapes):
            # Generate current shape path
            current_shape_path = f"{shape_path}/{shape_idx}" if shape_path else f"{shape_idx}"

            # Process shape based on type
            if isinstance(shape, GroupShape):
                # Process group shape and its children
                group_id = self._generate_id(f"group_{current_shape_path}_")

                # Create group element
                group_element = {
                    "element_id": group_id,
                    "doc_id": doc_id,
                    "element_type": "shape_group",
                    "parent_id": parent_id,
                    "content_preview": f"Shape Group {current_shape_path}",
                    "content_location": json.dumps({
                        "source": source_id,
                        "type": "shape_group",
                        "slide_index": slide_idx,
                        "shape_path": current_shape_path
                    }),
                    "content_hash": self._generate_hash(f"group_{current_shape_path}"),
                    "metadata": {
                        "slide_index": slide_idx,
                        "shape_index": shape_idx,
                        "shape_path": current_shape_path,
                        "shape_type": "group",
                        "shape_count": len(shape.shapes) if hasattr(shape, 'shapes') else 0
                    }
                }
                elements.append(group_element)

                # Process child shapes
                child_elements = self._process_shapes(shape.shapes, doc_id, group_id, source_id, slide_idx,
                                                      current_shape_path)
                elements.extend(child_elements)

            elif hasattr(shape, 'has_table') and shape.has_table and self.extract_tables:
                # Process table shape
                table_elements = self._process_table_shape(shape, doc_id, parent_id, source_id, slide_idx,
                                                           current_shape_path)
                elements.extend(table_elements)

            elif hasattr(shape, 'has_chart') and shape.has_chart and self.extract_charts:
                # Process chart shape
                chart_elements = self._process_chart_shape(shape, doc_id, parent_id, source_id, slide_idx,
                                                           current_shape_path)
                elements.extend(chart_elements)

            elif isinstance(shape, Picture) and self.extract_images:
                # Process picture shape
                picture_elements = self._process_picture_shape(shape, doc_id, parent_id, source_id, slide_idx,
                                                               current_shape_path)
                elements.extend(picture_elements)

            elif hasattr(shape, 'has_text_frame') and shape.has_text_frame and self.extract_text_boxes:
                # Process text shape
                text_elements = self._process_text_shape(shape, doc_id, parent_id, source_id, slide_idx,
                                                         current_shape_path)
                elements.extend(text_elements)

            elif self.extract_shapes:
                # Process generic shape
                shape_id = self._generate_id(f"shape_{current_shape_path}_")

                # Get shape name and type info
                shape_name = shape.name if hasattr(shape, 'name') else ""
                shape_type = self._get_shape_type(shape)

                # Create shape element
                shape_element = {
                    "element_id": shape_id,
                    "doc_id": doc_id,
                    "element_type": "shape",
                    "parent_id": parent_id,
                    "content_preview": f"Shape: {shape_name}" if shape_name else f"Shape {current_shape_path}",
                    "content_location": json.dumps({
                        "source": source_id,
                        "type": "shape",
                        "slide_index": slide_idx,
                        "shape_path": current_shape_path
                    }),
                    "content_hash": self._generate_hash(f"shape_{current_shape_path}"),
                    "metadata": {
                        "slide_index": slide_idx,
                        "shape_index": shape_idx,
                        "shape_path": current_shape_path,
                        "shape_type": shape_type,
                        "shape_name": shape_name
                    }
                }
                elements.append(shape_element)

        return elements

    def _process_text_shape(self, shape, doc_id: str, parent_id: str, source_id: str,
                            slide_idx: int, shape_path: str) -> List[Dict[str, Any]]:
        """
        Process a text shape into structured elements.

        Args:
            shape: The text shape
            doc_id: Document ID
            parent_id: Parent element ID
            source_id: Source identifier
            slide_idx: Slide index
            shape_path: Path to the shape

        Returns:
            List of text-related elements
        """
        elements = []

        try:
            if not shape.has_text_frame or not hasattr(shape, 'text_frame'):
                return elements

            text_frame = shape.text_frame
            text = text_frame.text

            if not text:
                return elements

            # Generate text element ID
            text_id = self._generate_id(f"text_{shape_path}_")

            # Get shape name if available
            shape_name = shape.name if hasattr(shape, 'name') else ""

            # Create text element
            text_element = {
                "element_id": text_id,
                "doc_id": doc_id,
                "element_type": "text_box",
                "parent_id": parent_id,
                "content_preview": text[:100] + ("..." if len(text) > 100 else ""),
                "content_location": json.dumps({
                    "source": source_id,
                    "type": "text_box",
                    "slide_index": slide_idx,
                    "shape_path": shape_path
                }),
                "content_hash": self._generate_hash(text),
                "metadata": {
                    "slide_index": slide_idx,
                    "shape_path": shape_path,
                    "shape_name": shape_name,
                    "text": text,
                    "is_title": self._is_title_shape(shape),
                    "level": self._get_paragraph_level(text_frame) if hasattr(text_frame, 'paragraphs') else 0
                }
            }
            elements.append(text_element)

            # Process paragraphs if detailed paragraph extraction is desired
            if hasattr(text_frame, 'paragraphs') and len(text_frame.paragraphs) > 1:
                for para_idx, paragraph in enumerate(text_frame.paragraphs):
                    para_text = paragraph.text
                    if not para_text:
                        continue

                    para_id = self._generate_id(f"para_{shape_path}_{para_idx}_")

                    para_element = {
                        "element_id": para_id,
                        "doc_id": doc_id,
                        "element_type": "paragraph",
                        "parent_id": text_id,
                        "content_preview": para_text[:100] + ("..." if len(para_text) > 100 else ""),
                        "content_location": json.dumps({
                            "source": source_id,
                            "type": "paragraph",
                            "slide_index": slide_idx,
                            "shape_path": shape_path,
                            "paragraph_index": para_idx
                        }),
                        "content_hash": self._generate_hash(para_text),
                        "metadata": {
                            "slide_index": slide_idx,
                            "shape_path": shape_path,
                            "paragraph_index": para_idx,
                            "text": para_text,
                            "level": paragraph.level if hasattr(paragraph, 'level') else 0
                        }
                    }
                    elements.append(para_element)

        except Exception as e:
            logger.warning(f"Error processing text shape: {str(e)}")

        return elements

    def _process_picture_shape(self, shape, doc_id: str, parent_id: str, source_id: str,
                               slide_idx: int, shape_path: str) -> List[Dict[str, Any]]:
        """
        Process a picture shape into structured elements.

        Args:
            shape: The picture shape
            doc_id: Document ID
            parent_id: Parent element ID
            source_id: Source identifier
            slide_idx: Slide index
            shape_path: Path to the shape

        Returns:
            List of picture-related elements
        """
        elements = []

        try:
            # Generate image element ID
            image_id = self._generate_id(f"image_{shape_path}_")

            # Get shape name and image info
            shape_name = shape.name if hasattr(shape, 'name') else ""
            image_name = ""

            if hasattr(shape, 'image') and hasattr(shape.image, 'filename'):
                image_name = shape.image.filename

            # Create image element
            image_element = {
                "element_id": image_id,
                "doc_id": doc_id,
                "element_type": "image",
                "parent_id": parent_id,
                "content_preview": f"Image: {image_name}" if image_name else f"Image {shape_path}",
                "content_location": json.dumps({
                    "source": source_id,
                    "type": "image",
                    "slide_index": slide_idx,
                    "shape_path": shape_path
                }),
                "content_hash": self._generate_hash(f"image_{shape_path}"),
                "metadata": {
                    "slide_index": slide_idx,
                    "shape_path": shape_path,
                    "shape_name": shape_name,
                    "image_name": image_name,
                    "alt_text": shape.alt_text if hasattr(shape, 'alt_text') else ""
                }
            }
            elements.append(image_element)

            # Process image caption (text) if available
            if hasattr(shape, 'text_frame') and shape.text_frame.text:
                caption_elements = self._process_text_shape(shape, doc_id, image_id, source_id, slide_idx,
                                                            f"{shape_path}_caption")
                elements.extend(caption_elements)

        except Exception as e:
            logger.warning(f"Error processing picture shape: {str(e)}")

        return elements

    def _process_table_shape(self, shape, doc_id: str, parent_id: str, source_id: str,
                             slide_idx: int, shape_path: str) -> List[Dict[str, Any]]:
        """
        Process a table shape into structured elements.

        Args:
            shape: The table shape
            doc_id: Document ID
            parent_id: Parent element ID
            source_id: Source identifier
            slide_idx: Slide index
            shape_path: Path to the shape

        Returns:
            List of table-related elements
        """
        elements = list()

        try:
            if not shape.has_table or not hasattr(shape, 'table'):
                return elements

            table = shape.table
            shape_name = shape.name if hasattr(shape, 'name') else ""

            # Generate table element ID
            table_id = self._generate_id(f"table_{shape_path}_")

            # Extract basic table content for preview
            table_text = ""
            for row in table.rows:
                for cell in row.cells:
                    if cell.text_frame.text:
                        table_text += cell.text_frame.text + " | "
                table_text += "\n"

            # Create table element
            table_element = {
                "element_id": table_id,
                "doc_id": doc_id,
                "element_type": "table",
                "parent_id": parent_id,
                "content_preview": f"Table: {table_text[:100]}" + ("..." if len(table_text) > 100 else ""),
                "content_location": json.dumps({
                    "source": source_id,
                    "type": "table",
                    "slide_index": slide_idx,
                    "shape_path": shape_path
                }),
                "content_hash": self._generate_hash(table_text),
                "metadata": {
                    "slide_index": slide_idx,
                    "shape_path": shape_path,
                    "shape_name": shape_name,
                    "rows": len(table.rows),
                    "columns": len(table.columns)
                }
            }
            elements.append(table_element)

            # Process table rows and cells
            for row_idx, row in enumerate(table.rows):
                # Generate row element ID
                row_id = self._generate_id(f"row_{shape_path}_{row_idx}_")

                # Create row element
                row_element = {
                    "element_id": row_id,
                    "doc_id": doc_id,
                    "element_type": "table_row",
                    "parent_id": table_id,
                    "content_preview": f"Row {row_idx + 1}",
                    "content_location": json.dumps({
                        "source": source_id,
                        "type": "table_row",
                        "slide_index": slide_idx,
                        "shape_path": shape_path,
                        "row": row_idx
                    }),
                    "content_hash": self._generate_hash(f"row_{shape_path}_{row_idx}"),
                    "metadata": {
                        "slide_index": slide_idx,
                        "shape_path": shape_path,
                        "row": row_idx
                    }
                }
                elements.append(row_element)

                # Process cells in this row
                for col_idx, cell in enumerate(row.cells):
                    cell_text = cell.text_frame.text if hasattr(cell, 'text_frame') else ""

                    if not cell_text:
                        continue

                    # Generate cell element ID
                    cell_id = self._generate_id(f"cell_{shape_path}_{row_idx}_{col_idx}_")

                    # Create cell element
                    cell_element = {
                        "element_id": cell_id,
                        "doc_id": doc_id,
                        "element_type": "table_cell",
                        "parent_id": row_id,
                        "content_preview": cell_text[:100] + ("..." if len(cell_text) > 100 else ""),
                        "content_location": json.dumps({
                            "source": source_id,
                            "type": "table_cell",
                            "slide_index": slide_idx,
                            "shape_path": shape_path,
                            "row": row_idx,
                            "col": col_idx
                        }),
                        "content_hash": self._generate_hash(cell_text),
                        "metadata": {
                            "slide_index": slide_idx,
                            "shape_path": shape_path,
                            "row": row_idx,
                            "col": col_idx,
                            "text": cell_text
                        }
                    }
                    elements.append(cell_element)

        except Exception as e:
            logger.warning(f"Error processing table shape: {str(e)}")

        return elements

    def _process_chart_shape(self, shape, doc_id: str, parent_id: str, source_id: str,
                             slide_idx: int, shape_path: str) -> List[Dict[str, Any]]:
        """
        Process a chart shape into structured elements.

        Args:
            shape: The chart shape
            doc_id: Document ID
            parent_id: Parent element ID
            source_id: Source identifier
            slide_idx: Slide index
            shape_path: Path to the shape

        Returns:
            List of chart-related elements
        """
        elements = []

        try:
            if not shape.has_chart or not hasattr(shape, 'chart'):
                return elements

            chart = shape.chart
            shape_name = shape.name if hasattr(shape, 'name') else ""

            # Generate chart element ID
            chart_id = self._generate_id(f"chart_{shape_path}_")

            # Get chart type
            chart_type = "unknown"
            if hasattr(chart, 'chart_type'):
                chart_type = str(chart.chart_type)

            # Get chart title
            chart_title = ""
            if hasattr(chart, 'chart_title') and hasattr(chart.chart_title, 'text_frame'):
                chart_title = chart.chart_title.text_frame.text

            # Create chart element
            chart_element = {
                "element_id": chart_id,
                "doc_id": doc_id,
                "element_type": "chart",
                "parent_id": parent_id,
                "content_preview": f"Chart: {chart_title}" if chart_title else f"Chart {shape_path}",
                "content_location": json.dumps({
                    "source": source_id,
                    "type": "chart",
                    "slide_index": slide_idx,
                    "shape_path": shape_path
                }),
                "content_hash": self._generate_hash(f"chart_{shape_path}"),
                "metadata": {
                    "slide_index": slide_idx,
                    "shape_path": shape_path,
                    "shape_name": shape_name,
                    "chart_type": chart_type,
                    "chart_title": chart_title
                }
            }
            elements.append(chart_element)

            # Extract chart category and series names if available
            if hasattr(chart, 'plots') and chart.plots:
                plot = chart.plots[0]

                if hasattr(plot, 'categories'):
                    categories = []
                    for category in plot.categories:
                        if category:
                            categories.append(str(category))

                    if categories:
                        chart_element["metadata"]["categories"] = categories

                if hasattr(plot, 'series'):
                    series_names = []
                    for series in plot.series:
                        if hasattr(series, 'name') and series.name:
                            series_names.append(str(series.name))

                    if series_names:
                        chart_element["metadata"]["series"] = series_names

        except Exception as e:
            logger.warning(f"Error processing chart shape: {str(e)}")

        return elements

    def _process_notes(self, notes_slide, slide_idx: int, doc_id: str, parent_id: str, source_id: str) -> List[
        Dict[str, Any]]:
        """
        Process slide notes into structured elements.

        Args:
            notes_slide: The notes slide
            slide_idx: Slide index
            doc_id: Document ID
            parent_id: Parent element ID
            source_id: Source identifier

        Returns:
            List of notes-related elements
        """
        elements = []

        try:
            if not notes_slide or not hasattr(notes_slide, 'notes_text_frame'):
                return elements

            notes_text = notes_slide.notes_text_frame.text

            if not notes_text:
                return elements

            # Generate notes element ID
            notes_id = self._generate_id(f"notes_{slide_idx}_")

            # Create notes element
            notes_element = {
                "element_id": notes_id,
                "doc_id": doc_id,
                "element_type": "slide_notes",
                "parent_id": parent_id,
                "content_preview": notes_text[:100] + ("..." if len(notes_text) > 100 else ""),
                "content_location": json.dumps({
                    "source": source_id,
                    "type": "slide_notes",
                    "slide_index": slide_idx
                }),
                "content_hash": self._generate_hash(notes_text),
                "metadata": {
                    "slide_index": slide_idx,
                    "text": notes_text
                }
            }
            elements.append(notes_element)

        except Exception as e:
            logger.warning(f"Error processing slide notes: {str(e)}")

        return elements

    def _process_comments(self, slide, slide_idx: int, doc_id: str, parent_id: str, source_id: str) -> List[
        Dict[str, Any]]:
        """
        Process slide comments into structured elements.

        Args:
            slide: The slide
            slide_idx: Slide index
            doc_id: Document ID
            parent_id: Parent element ID
            source_id: Source identifier

        Returns:
            List of comment-related elements
        """
        elements = []

        try:
            # Check if slide has comments
            if not hasattr(slide, 'comments') or not slide.comments:
                return elements

            # Create comments container element
            comments_id = self._generate_id(f"comments_{slide_idx}_")

            comments_element = {
                "element_id": comments_id,
                "doc_id": doc_id,
                "element_type": "comments_container",
                "parent_id": parent_id,
                "content_preview": f"Comments for Slide {slide_idx + 1}",
                "content_location": json.dumps({
                    "source": source_id,
                    "type": "comments_container",
                    "slide_index": slide_idx
                }),
                "content_hash": self._generate_hash(f"comments_{slide_idx}"),
                "metadata": {
                    "slide_index": slide_idx,
                    "comment_count": len(slide.comments)
                }
            }
            elements.append(comments_element)

            # Process individual comments
            for comment_idx, comment in enumerate(slide.comments):
                comment_text = comment.text if hasattr(comment, 'text') else ""

                if not comment_text:
                    continue

                # Generate comment element ID
                comment_id = self._generate_id(f"comment_{slide_idx}_{comment_idx}_")

                # Get comment author and date if available
                author = comment.author if hasattr(comment, 'author') else "Unknown"
                date = comment.date if hasattr(comment, 'date') else None

                # Create comment element
                comment_element = {
                    "element_id": comment_id,
                    "doc_id": doc_id,
                    "element_type": "comment",
                    "parent_id": comments_id,
                    "content_preview": comment_text[:100] + ("..." if len(comment_text) > 100 else ""),
                    "content_location": json.dumps({
                        "source": source_id,
                        "type": "comment",
                        "slide_index": slide_idx,
                        "comment_index": comment_idx
                    }),
                    "content_hash": self._generate_hash(comment_text),
                    "metadata": {
                        "slide_index": slide_idx,
                        "comment_index": comment_idx,
                        "author": author,
                        "date": date.timestamp() if date and hasattr(date, 'timestamp') else None,
                        "text": comment_text
                    }
                }
                elements.append(comment_element)

        except Exception as e:
            logger.warning(f"Error processing slide comments: {str(e)}")

        return elements

    def _extract_slide_masters(self, presentation, doc_id: str, parent_id: str, source_id: str) -> List[Dict[str, Any]]:
        """
        Extract slide masters from presentation.

        Args:
            presentation: The PPTX presentation
            doc_id: Document ID
            parent_id: Parent element ID
            source_id: Source identifier

        Returns:
            List of slide master elements
        """
        elements = []

        try:
            if not hasattr(presentation, 'slide_masters'):
                return elements

            # Create masters container element
            masters_id = self._generate_id("masters_")

            masters_element = {
                "element_id": masters_id,
                "doc_id": doc_id,
                "element_type": "slide_masters",
                "parent_id": parent_id,
                "content_preview": "Slide Masters",
                "content_location": json.dumps({
                    "source": source_id,
                    "type": "slide_masters"
                }),
                "content_hash": "",
                "metadata": {
                    "master_count": len(presentation.slide_masters)
                }
            }
            elements.append(masters_element)

            # Process individual masters
            for master_idx, master in enumerate(presentation.slide_masters):
                # Generate master ID
                master_id = self._generate_id(f"master_{master_idx}_")

                # Create master element
                master_element = {
                    "element_id": master_id,
                    "doc_id": doc_id,
                    "element_type": "slide_master",
                    "parent_id": masters_id,
                    "content_preview": f"Slide Master {master_idx + 1}",
                    "content_location": json.dumps({
                        "source": source_id,
                        "type": "slide_master",
                        "index": master_idx
                    }),
                    "content_hash": self._generate_hash(f"master_{master_idx}"),
                    "metadata": {
                        "index": master_idx,
                        "layout_count": len(master.slide_layouts) if hasattr(master, 'slide_layouts') else 0
                    }
                }
                elements.append(master_element)

                # Process master shapes if desired
                if self.extract_shapes and hasattr(master, 'shapes'):
                    shape_elements = self._process_shapes(master.shapes, doc_id, master_id, source_id, -1,
                                                          f"master_{master_idx}")
                    elements.extend(shape_elements)

        except Exception as e:
            logger.warning(f"Error extracting slide masters: {str(e)}")

        return elements

    def _extract_slide_templates(self, presentation, doc_id: str, parent_id: str, source_id: str) -> List[
        Dict[str, Any]]:
        """
        Extract slide templates (layouts) from presentation.

        Args:
            presentation: The PPTX presentation
            doc_id: Document ID
            parent_id: Parent element ID
            source_id: Source identifier

        Returns:
            List of slide template elements
        """
        elements = []

        try:
            # Collect all slide layouts from all masters
            layouts = []
            layout_names = set()

            if hasattr(presentation, 'slide_masters'):
                for master in presentation.slide_masters:
                    if hasattr(master, 'slide_layouts'):
                        for layout in master.slide_layouts:
                            # Avoid duplicates by name
                            layout_name = layout.name if hasattr(layout, 'name') else ""
                            if layout_name not in layout_names:
                                layouts.append(layout)
                                layout_names.add(layout_name)

            if not layouts:
                return elements

            # Create templates container element
            templates_id = self._generate_id("templates_")

            templates_element = {
                "element_id": templates_id,
                "doc_id": doc_id,
                "element_type": "slide_templates",
                "parent_id": parent_id,
                "content_preview": "Slide Templates",
                "content_location": json.dumps({
                    "source": source_id,
                    "type": "slide_templates"
                }),
                "content_hash": "",
                "metadata": {
                    "template_count": len(layouts)
                }
            }
            elements.append(templates_element)

            # Process individual templates
            for layout_idx, layout in enumerate(layouts):
                # Generate layout ID
                layout_id = self._generate_id(f"layout_{layout_idx}_")

                # Get layout name
                layout_name = layout.name if hasattr(layout, 'name') else f"Layout {layout_idx + 1}"

                # Create layout element
                layout_element = {
                    "element_id": layout_id,
                    "doc_id": doc_id,
                    "element_type": "slide_layout",
                    "parent_id": templates_id,
                    "content_preview": layout_name,
                    "content_location": json.dumps({
                        "source": source_id,
                        "type": "slide_layout",
                        "index": layout_idx
                    }),
                    "content_hash": self._generate_hash(f"layout_{layout_idx}"),
                    "metadata": {
                        "index": layout_idx,
                        "name": layout_name
                    }
                }
                elements.append(layout_element)

                # Process layout shapes if desired
                if self.extract_shapes and hasattr(layout, 'shapes'):
                    shape_elements = self._process_shapes(layout.shapes, doc_id, layout_id, source_id, -1,
                                                          f"layout_{layout_idx}")
                    elements.extend(shape_elements)

        except Exception as e:
            logger.warning(f"Error extracting slide templates: {str(e)}")

        return elements

    @staticmethod
    def _get_slide_title(slide: Slide) -> str:
        """Get slide title text."""
        try:
            # Look for a shape with placeholder type as title
            for shape in slide.shapes:
                if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
                    if hasattr(shape, 'placeholder_format') and shape.placeholder_format.type == 1:  # 1 = TITLE
                        if hasattr(shape, 'text_frame') and shape.text_frame.text:
                            return shape.text_frame.text

            # If no title placeholder, look for first shape with text
            for shape in slide.shapes:
                if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                    if shape.text_frame.text:
                        return shape.text_frame.text

        except Exception as e:
            logger.debug(f"Error getting slide title: {str(e)}")

        return ""

    @staticmethod
    def _get_slide_layout_name(slide: Slide) -> str:
        """Get slide layout name."""
        try:
            if hasattr(slide, 'slide_layout') and hasattr(slide.slide_layout, 'name'):
                return slide.slide_layout.name
        except Exception as e:
            logger.debug(f"Error getting slide layout name: {str(e)}")

        return "Unknown Layout"

    @staticmethod
    def _is_title_shape(shape) -> bool:
        """Check if shape is a title placeholder."""
        try:
            if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
                if hasattr(shape, 'placeholder_format') and shape.placeholder_format.type in [1,
                                                                                              2]:  # 1 = TITLE, 2 = CENTERED_TITLE
                    return True
        except Exception:
            pass

        return False

    @staticmethod
    def _get_paragraph_level(text_frame: TextFrame) -> int:
        """Get the outline level of text frame."""
        try:
            if hasattr(text_frame, 'paragraphs') and text_frame.paragraphs:
                return text_frame.paragraphs[0].level if hasattr(text_frame.paragraphs[0], 'level') else 0
        except Exception:
            pass

        return 0

    @staticmethod
    def _get_shape_type(shape) -> str:
        """Get the type of shape."""
        try:
            if hasattr(shape, 'shape_type'):
                return str(shape.shape_type)

            if isinstance(shape, Picture):
                return "picture"

            if hasattr(shape, 'has_table') and shape.has_table:
                return "table"

            if hasattr(shape, 'has_chart') and shape.has_chart:
                return "chart"

            if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                return "text"

            if isinstance(shape, GroupShape):
                return "group"

        except Exception:
            pass

        return "shape"

    @staticmethod
    def _generate_hash(content: str) -> str:
        """
        Generate a hash of content for change detection.

        Args:
            content: Text content

        Returns:
            MD5 hash of content
        """
        import hashlib
        return hashlib.md5(content.encode('utf-8')).hexdigest()
